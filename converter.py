import asyncio
import os
import zipfile

import pdfplumber
import docx
from pptx import Presentation
from openpyxl import load_workbook
import aiofiles
from io import BytesIO

ALLOWED = {'txt', 'pdf', 'docx', 'pptx', 'xlsx'}

def create_zip_directory(directory_path, family_dir):
    memory_file = BytesIO()
    with zipfile.ZipFile(memory_file, 'w') as zf:
        for root, dirs, files in os.walk(directory_path):
            for file in files:
                if file.endswith('.txt'):
                    pdf_filename = file.replace('.txt', '.pdf')
                    pdf_path = os.path.join(family_dir, pdf_filename)
                    if os.path.exists(pdf_path):
                        zf.write(pdf_path, os.path.basename(pdf_path))
                    else:
                        print(f"Предупреждение: файл {pdf_filename} не найден в {family_dir}")

    memory_file.seek(0)
    return memory_file



async def convert_pdf_to_txt(pdf_file, output_folder):
    try:
        txt_filename = os.path.splitext(pdf_file.filename)[0] + '.txt'
        txt_path = os.path.join(output_folder, txt_filename)

        with pdfplumber.open(pdf_file) as pdf:
            async with aiofiles.open(txt_path, 'w', encoding='utf-8') as txt_file:
                for page in pdf.pages:
                    text = page.extract_text()
                    if text:
                        await txt_file.write(text + '\n')

        return txt_filename
    except Exception as e:
        print(f"Error converting PDF to TXT: {e}")
        return None

async def convert_docx_to_txt(docx_file, output_folder):
    try:
        txt_filename = os.path.splitext(docx_file.filename)[0] + '.txt'
        txt_path = os.path.join(output_folder, txt_filename)

        doc = docx.Document(docx_file)
        async with aiofiles.open(txt_path, 'w', encoding='utf-8') as txt_file:
            for para in doc.paragraphs:
                await txt_file.write(para.text + '\n')

        return txt_filename
    except Exception as e:
        print(f"Error converting DOCX to TXT: {e}")
        return None

async def convert_pptx_to_txt(pptx_file, output_folder):
    try:
        txt_filename = os.path.splitext(pptx_file.filename)[0] + '.txt'
        txt_path = os.path.join(output_folder, txt_filename)

        presentation = Presentation(pptx_file)
        async with aiofiles.open(txt_path, 'w', encoding='utf-8') as txt_file:
            for slide in presentation.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        await txt_file.write(shape.text + '\n')

        return txt_filename
    except Exception as e:
        print(f"Error converting PPTX to TXT: {e}")
        return None

async def convert_xlsx_to_txt(xlsx_file, output_folder):
    try:
        txt_filename = os.path.splitext(xlsx_file.filename)[0] + '.txt'
        txt_path = os.path.join(output_folder, txt_filename)

        workbook = load_workbook(xlsx_file, data_only=True)
        async with aiofiles.open(txt_path, 'w', encoding='utf-8') as txt_file:
            for sheet in workbook:
                await txt_file.write(f"Sheet: {sheet.title}\n")
                for row in sheet.iter_rows(values_only=True):
                    row_data = '\t'.join([str(cell) for cell in row])
                    await txt_file.write(row_data + '\n')

        return txt_filename
    except Exception as e:
        print(f"Error converting XLSX to TXT: {e}")
        return None

async def convert_to_pdf(input_file, output_folder):
    try:
        if not os.path.exists(input_file):
            raise FileNotFoundError(f"Файл {input_file} не найден.")

        pdf_filename = os.path.splitext(os.path.basename(input_file))[0] + '.pdf'
        pdf_output_path = os.path.join(output_folder, pdf_filename)

        #LibreOffice
        process = await asyncio.create_subprocess_exec(
            'libreoffice', '--headless', '--convert-to', 'pdf', '--outdir',
            output_folder, input_file,
            stdout=asyncio.subprocess.PIPE, stderr=asyncio.subprocess.PIPE
        )
        stdout, stderr = await process.communicate()

        if process.returncode == 0:
            print(f"Файл успешно конвертирован: {pdf_output_path}")
            return pdf_output_path
        else:
            print(f"Ошибка при конвертации: {stderr.decode()}")
            return None

    except Exception as e:
        print(f"Ошибка: {e}")
        return None

def allowed_file(filename):
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in ALLOWED
